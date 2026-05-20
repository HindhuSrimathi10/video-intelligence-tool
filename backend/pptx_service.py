from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import datetime
import re

# Color Scheme
DARK_BG = RGBColor(15, 23, 42)
ACCENT_BLUE = RGBColor(59, 130, 246)
ACCENT_GREEN = RGBColor(16, 185, 129)
ACCENT_ORANGE = RGBColor(245, 158, 11)
ACCENT_RED = RGBColor(239, 68, 68)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(148, 163, 184)
CARD_BG = RGBColor(30, 41, 59)

COMPANY_COLORS = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED, RGBColor(168, 85, 247)]

def set_slide_background(slide, prs, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, left, top, width, height, font_size=18,
                 bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_rectangle(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def parse_insights_sections(insights_text):
    sections = {
        "executive_summary": "",
        "channel_overview": "",
        "content_performance": "",
        "content_topics": "",
        "posting_frequency": "",
        "engagement_analysis": "",
        "gap_analysis": "",
        "recommendations": "",
        "rankings": ""
    }
    patterns = {
        "executive_summary": r"1\.\s*EXECUTIVE SUMMARY(.*?)(?=2\.|$)",
        "channel_overview": r"2\.\s*CHANNEL OVERVIEW[^\n]*(.*?)(?=3\.|$)",
        "content_performance": r"3\.\s*CONTENT PERFORMANCE[^\n]*(.*?)(?=4\.|$)",
        "content_topics": r"4\.\s*CONTENT TOPICS[^\n]*(.*?)(?=5\.|$)",
        "posting_frequency": r"5\.\s*POSTING FREQUENCY[^\n]*(.*?)(?=6\.|$)",
        "engagement_analysis": r"6\.\s*ENGAGEMENT ANALYSIS[^\n]*(.*?)(?=7\.|$)",
        "gap_analysis": r"7\.\s*GAP ANALYSIS[^\n]*(.*?)(?=8\.|$)",
        "recommendations": r"8\.\s*VIDEO MARKETING RECOMMENDATIONS[^\n]*(.*?)(?=9\.|$)",
        "rankings": r"9\.\s*COMPANY RANKINGS[^\n]*(.*?)(?=$)"
    }
    for key, pattern in patterns.items():
        match = re.search(pattern, insights_text, re.DOTALL | re.IGNORECASE)
        if match:
            sections[key] = match.group(1).strip()[:800]
    return sections

def create_cover_slide(prs, companies, main_company):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    # Top accent bar
    add_rectangle(slide, 0, 0, 10, 0.08, ACCENT_BLUE)

    # Main title
    add_text_box(slide, "VIDEO COMPETITOR", 0.5, 1.2, 9, 1.2,
                font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "INTELLIGENCE REPORT", 0.5, 2.2, 9, 1.2,
                font_size=44, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # Divider line
    add_rectangle(slide, 3, 3.3, 4, 0.05, ACCENT_GREEN)

    # Company name
    add_text_box(slide, f"Primary Company: {main_company}", 0.5, 3.6, 9, 0.6,
                font_size=20, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    # Competitors
    competitors = [c for c in companies if c != main_company]
    comp_text = "Competitors: " + " | ".join(competitors)
    add_text_box(slide, comp_text, 0.5, 4.3, 9, 0.6,
                font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Date
    date_str = datetime.datetime.now().strftime("%B %d, %Y")
    add_text_box(slide, f"Report Generated: {date_str}", 0.5, 5.2, 9, 0.5,
                font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Bottom bar
    add_rectangle(slide, 0, 6.9, 10, 0.1, ACCENT_BLUE)
    add_text_box(slide, "CONFIDENTIAL — Video Marketing Intelligence", 0.5, 6.95, 9, 0.3,
                font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

def create_executive_summary_slide(prs, insights_sections, all_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    add_rectangle(slide, 0, 0, 10, 0.08, ACCENT_GREEN)
    add_text_box(slide, "02", 0.3, 0.15, 1, 0.5, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, "EXECUTIVE SUMMARY", 0.5, 0.15, 9, 0.6,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rectangle(slide, 0.5, 0.85, 9, 0.04, ACCENT_GREEN)

    summary_text = insights_sections.get("executive_summary", "Analysis complete.")
    summary_text = summary_text[:600] if len(summary_text) > 600 else summary_text

    add_rectangle(slide, 0.4, 1.0, 9.2, 3.5, CARD_BG)
    add_text_box(slide, summary_text, 0.7, 1.1, 8.7, 3.3,
                font_size=13, color=WHITE)

    # Quick stats row
    companies = list(all_data.keys())
    col_width = 9.0 / len(companies)
    for i, (company, data) in enumerate(all_data.items()):
        x = 0.5 + i * col_width
        color = COMPANY_COLORS[i % len(COMPANY_COLORS)]
        add_rectangle(slide, x, 4.7, col_width - 0.1, 1.5, CARD_BG)
        add_rectangle(slide, x, 4.7, col_width - 0.1, 0.08, color)
        subs = data.get('channel_stats', {}).get('subscriber_count', 0) if data.get('channel_stats') else 0
        add_text_box(slide, company[:12], x + 0.05, 4.8, col_width - 0.2, 0.4,
                    font_size=10, bold=True, color=color)
        add_text_box(slide, f"{subs:,}", x + 0.05, 5.2, col_width - 0.2, 0.4,
                    font_size=14, bold=True, color=WHITE)
        add_text_box(slide, "subscribers", x + 0.05, 5.6, col_width - 0.2, 0.3,
                    font_size=8, color=LIGHT_GRAY)

def create_channel_overview_slide(prs, all_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    add_rectangle(slide, 0, 0, 10, 0.08, ACCENT_BLUE)
    add_text_box(slide, "03", 0.3, 0.15, 1, 0.5, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, "CHANNEL OVERVIEW COMPARISON", 0.5, 0.15, 9, 0.6,
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rectangle(slide, 0.5, 0.85, 9, 0.04, ACCENT_BLUE)

    companies = list(all_data.keys())
    col_width = 9.0 / len(companies)

    for i, (company, data) in enumerate(all_data.items()):
        x = 0.5 + i * col_width
        color = COMPANY_COLORS[i % len(COMPANY_COLORS)]
        stats = data.get('channel_stats', {}) or {}

        add_rectangle(slide, x, 1.0, col_width - 0.1, 5.8, CARD_BG)
        add_rectangle(slide, x, 1.0, col_width - 0.1, 0.1, color)

        add_text_box(slide, company[:14], x + 0.1, 1.15, col_width - 0.2, 0.5,
                    font_size=13, bold=True, color=color)

        metrics = [
            ("SUBSCRIBERS", f"{stats.get('subscriber_count', 0):,}"),
            ("TOTAL VIDEOS", f"{stats.get('total_videos', 0):,}"),
            ("TOTAL VIEWS", f"{stats.get('total_views', 0):,}"),
            ("UPLOAD FREQ", data.get('upload_frequency', 'N/A')[:20]),
            ("AVG VIEWS", f"{data.get('avg_views', 0):,}"),
            ("AVG LIKES", f"{data.get('avg_likes', 0):,}"),
            ("AVG COMMENTS", f"{data.get('avg_comments', 0):,}"),
        ]

        for j, (label, value) in enumerate(metrics):
            y = 1.8 + j * 0.7
            add_text_box(slide, label, x + 0.1, y, col_width - 0.2, 0.25,
                        font_size=7, color=LIGHT_GRAY, bold=True)
            add_text_box(slide, str(value), x + 0.1, y + 0.25, col_width - 0.2, 0.35,
                        font_size=11, color=WHITE, bold=True)

def create_content_performance_slide(prs, all_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    add_rectangle(slide, 0, 0, 10, 0.08, ACCENT_ORANGE)
    add_text_box(slide, "04", 0.3, 0.15, 1, 0.5, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, "CONTENT PERFORMANCE — TOP VIDEOS", 0.5, 0.15, 9, 0.6,
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rectangle(slide, 0.5, 0.85, 9, 0.04, ACCENT_ORANGE)

    companies = list(all_data.keys())
    col_width = 9.0 / len(companies)

    for i, (company, data) in enumerate(all_data.items()):
        x = 0.5 + i * col_width
        color = COMPANY_COLORS[i % len(COMPANY_COLORS)]
        videos = data.get('videos', [])

        add_rectangle(slide, x, 1.0, col_width - 0.1, 5.8, CARD_BG)
        add_rectangle(slide, x, 1.0, col_width - 0.1, 0.1, color)
        add_text_box(slide, company[:14], x + 0.1, 1.15, col_width - 0.2, 0.4,
                    font_size=12, bold=True, color=color)

        if videos:
            for j, video in enumerate(videos[:3]):
                y = 1.7 + j * 1.7
                add_rectangle(slide, x + 0.05, y, col_width - 0.2, 1.5,
                             RGBColor(40, 55, 75))
                title = video['title'][:45] + "..." if len(video['title']) > 45 else video['title']
                add_text_box(slide, title, x + 0.1, y + 0.05, col_width - 0.25, 0.7,
                            font_size=8, color=WHITE)
                add_text_box(slide, f"👁 {video['view_count']:,}",
                            x + 0.1, y + 0.75, col_width - 0.25, 0.3,
                            font_size=8, color=ACCENT_BLUE)
                add_text_box(slide, f"👍 {video['like_count']:,}",
                            x + 0.1, y + 1.05, col_width - 0.25, 0.3,
                            font_size=8, color=ACCENT_GREEN)
        else:
            add_text_box(slide, "No video data available", x + 0.1, 2.0,
                        col_width - 0.2, 0.5, font_size=10, color=LIGHT_GRAY)

def create_engagement_chart_slide(prs, all_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    add_rectangle(slide, 0, 0, 10, 0.08, ACCENT_GREEN)
    add_text_box(slide, "05", 0.3, 0.15, 1, 0.5, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, "ENGAGEMENT ANALYSIS", 0.5, 0.15, 9, 0.6,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rectangle(slide, 0.5, 0.85, 9, 0.04, ACCENT_GREEN)

    companies = list(all_data.keys())
    col_width = 9.0 / len(companies)

    for i, (company, data) in enumerate(all_data.items()):
        x = 0.5 + i * col_width
        color = COMPANY_COLORS[i % len(COMPANY_COLORS)]

        add_rectangle(slide, x, 1.0, col_width - 0.1, 5.8, CARD_BG)
        add_rectangle(slide, x, 1.0, col_width - 0.1, 0.1, color)
        add_text_box(slide, company[:14], x + 0.1, 1.15, col_width - 0.2, 0.4,
                    font_size=12, bold=True, color=color)

        metrics = [
            ("AVG VIEWS", f"{data.get('avg_views', 0):,}", ACCENT_BLUE),
            ("AVG LIKES", f"{data.get('avg_likes', 0):,}", ACCENT_GREEN),
            ("AVG COMMENTS", f"{data.get('avg_comments', 0):,}", ACCENT_ORANGE),
        ]

        for j, (label, value, mcolor) in enumerate(metrics):
            y = 1.8 + j * 1.6
            add_rectangle(slide, x + 0.1, y, col_width - 0.25, 1.4, RGBColor(40, 55, 75))
            add_text_box(slide, label, x + 0.15, y + 0.1, col_width - 0.35, 0.3,
                        font_size=8, color=LIGHT_GRAY, bold=True)
            add_text_box(slide, value, x + 0.15, y + 0.45, col_width - 0.35, 0.6,
                        font_size=18, bold=True, color=mcolor)

        engagement_text = insights_text = ""
        top_video = data.get("top_video")
        if top_video:
            engagement_text = f"Top: {top_video['view_count']:,} views"
        add_text_box(slide, engagement_text, x + 0.1, 6.6, col_width - 0.2, 0.3,
                    font_size=8, color=LIGHT_GRAY)

def create_topics_themes_slide(prs, all_data, insights_sections):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    add_rectangle(slide, 0, 0, 10, 0.08, ACCENT_BLUE)
    add_text_box(slide, "06", 0.3, 0.15, 1, 0.5, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, "CONTENT TOPICS & THEMES", 0.5, 0.15, 9, 0.6,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rectangle(slide, 0.5, 0.85, 9, 0.04, ACCENT_BLUE)

    topics_text = insights_sections.get("content_topics", "")
    add_rectangle(slide, 0.4, 1.0, 9.2, 5.8, CARD_BG)
    add_text_box(slide, topics_text[:700] if topics_text else "Content topic analysis based on video titles and descriptions.",
                0.7, 1.1, 8.7, 5.5, font_size=12, color=WHITE)

def create_posting_frequency_slide(prs, all_data, insights_sections):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    add_rectangle(slide, 0, 0, 10, 0.08, ACCENT_ORANGE)
    add_text_box(slide, "07", 0.3, 0.15, 1, 0.5, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, "POSTING FREQUENCY & CONSISTENCY", 0.5, 0.15, 9, 0.6,
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rectangle(slide, 0.5, 0.85, 9, 0.04, ACCENT_ORANGE)

    companies = list(all_data.keys())
    col_width = 9.0 / len(companies)

    for i, (company, data) in enumerate(all_data.items()):
        x = 0.5 + i * col_width
        color = COMPANY_COLORS[i % len(COMPANY_COLORS)]

        add_rectangle(slide, x, 1.0, col_width - 0.1, 3.0, CARD_BG)
        add_rectangle(slide, x, 1.0, col_width - 0.1, 0.1, color)
        add_text_box(slide, company[:14], x + 0.1, 1.15, col_width - 0.2, 0.4,
                    font_size=12, bold=True, color=color)

        freq = data.get('upload_frequency', 'N/A')
        add_text_box(slide, freq[:30], x + 0.1, 1.7, col_width - 0.2, 0.6,
                    font_size=10, color=WHITE)

        stats = data.get('channel_stats', {}) or {}
        total_videos = stats.get('total_videos', 0)
        add_text_box(slide, f"Total: {total_videos:,} videos",
                    x + 0.1, 2.4, col_width - 0.2, 0.4,
                    font_size=10, color=LIGHT_GRAY)

    freq_text = insights_sections.get("posting_frequency", "")
    add_rectangle(slide, 0.4, 4.2, 9.2, 2.5, CARD_BG)
    add_text_box(slide, "AI Analysis:", 0.7, 4.3, 9, 0.4,
                font_size=11, bold=True, color=ACCENT_ORANGE)
    add_text_box(slide, freq_text[:400] if freq_text else "Posting frequency analysis.",
                0.7, 4.7, 8.7, 2.0, font_size=11, color=WHITE)

def create_gap_analysis_slide(prs, insights_sections):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    add_rectangle(slide, 0, 0, 10, 0.08, ACCENT_RED)
    add_text_box(slide, "08", 0.3, 0.15, 1, 0.5, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, "GAP ANALYSIS — MISSED OPPORTUNITIES", 0.5, 0.15, 9, 0.6,
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rectangle(slide, 0.5, 0.85, 9, 0.04, ACCENT_RED)

    gap_text = insights_sections.get("gap_analysis", "")
    add_rectangle(slide, 0.4, 1.0, 9.2, 5.8, CARD_BG)
    add_text_box(slide, "🔍 Content & Strategy Gaps Identified:", 0.7, 1.1, 9, 0.5,
                font_size=14, bold=True, color=ACCENT_RED)
    add_text_box(slide, gap_text[:700] if gap_text else "Gap analysis based on competitive data.",
                0.7, 1.7, 8.7, 5.0, font_size=12, color=WHITE)

def create_recommendations_slide(prs, insights_sections, main_company):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    add_rectangle(slide, 0, 0, 10, 0.08, ACCENT_GREEN)
    add_text_box(slide, "09", 0.3, 0.15, 1, 0.5, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, f"RECOMMENDATIONS FOR {main_company.upper()}", 0.5, 0.15, 9, 0.6,
                font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rectangle(slide, 0.5, 0.85, 9, 0.04, ACCENT_GREEN)

    rec_text = insights_sections.get("recommendations", "")
    add_rectangle(slide, 0.4, 1.0, 9.2, 5.8, CARD_BG)
    add_text_box(slide, "🚀 Strategic Action Items:", 0.7, 1.1, 9, 0.5,
                font_size=14, bold=True, color=ACCENT_GREEN)
    add_text_box(slide, rec_text[:700] if rec_text else "Strategic recommendations based on competitive analysis.",
                0.7, 1.7, 8.7, 5.0, font_size=12, color=WHITE)

def create_rankings_slide(prs, all_data, insights_sections):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    add_rectangle(slide, 0, 0, 10, 0.08, ACCENT_BLUE)
    add_text_box(slide, "10", 0.3, 0.15, 1, 0.5, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, "COMPANY RANKINGS & SCORES", 0.5, 0.15, 9, 0.6,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rectangle(slide, 0.5, 0.85, 9, 0.04, ACCENT_BLUE)

    companies = list(all_data.keys())
    scores = []
    for company, data in all_data.items():
        score = 0
        stats = data.get('channel_stats', {}) or {}
        subs = stats.get('subscriber_count', 0)
        videos = stats.get('total_videos', 0)
        avg_views = data.get('avg_views', 0)

        if subs > 1000000: score += 3
        elif subs > 100000: score += 2
        elif subs > 10000: score += 1

        if videos > 500: score += 2
        elif videos > 100: score += 1

        if avg_views > 100000: score += 3
        elif avg_views > 10000: score += 2
        elif avg_views > 1000: score += 1

        freq = data.get('upload_frequency', '')
        if 'Very Active' in freq: score += 2
        elif 'Weekly' in freq: score += 1

        scores.append((company, min(score, 10)))

    scores.sort(key=lambda x: x[1], reverse=True)

    for i, (company, score) in enumerate(scores):
        y = 1.1 + i * 1.0
        color = COMPANY_COLORS[list(all_data.keys()).index(company) % len(COMPANY_COLORS)]

        medal = ["🥇", "🥈", "🥉", "4️⃣", "5️⃣"][i] if i < 5 else str(i+1)
        add_rectangle(slide, 0.4, y, 9.2, 0.85, CARD_BG)
        add_rectangle(slide, 0.4, y, 0.08, 0.85, color)

        add_text_box(slide, medal, 0.6, y + 0.1, 0.6, 0.6, font_size=18)
        add_text_box(slide, company, 1.3, y + 0.1, 4, 0.6,
                    font_size=16, bold=True, color=color)

        bar_width = (score / 10) * 3.5
        add_rectangle(slide, 5.5, y + 0.3, bar_width, 0.3, color)
        add_text_box(slide, f"{score}/10", 9.1, y + 0.2, 0.8, 0.5,
                    font_size=14, bold=True, color=WHITE)

def create_summary_slide(prs, all_data, main_company):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs, DARK_BG)

    add_rectangle(slide, 0, 0, 10, 0.08, ACCENT_BLUE)
    add_text_box(slide, "THANK YOU", 0.5, 1.5, 9, 1.2,
                font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Video Intelligence Report Complete", 0.5, 2.8, 9, 0.8,
                font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_rectangle(slide, 3, 3.8, 4, 0.05, ACCENT_GREEN)

    add_text_box(slide, f"Prepared for: {main_company}", 0.5, 4.1, 9, 0.5,
                font_size=14, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    date_str = datetime.datetime.now().strftime("%B %d, %Y")
    add_text_box(slide, f"Generated on {date_str}", 0.5, 4.7, 9, 0.5,
                font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_rectangle(slide, 0, 6.9, 10, 0.1, ACCENT_BLUE)

def create_pptx_report(all_data, insights, main_company):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    companies = list(all_data.keys())
    insights_sections = parse_insights_sections(insights)

    # Create all slides
    create_cover_slide(prs, companies, main_company)
    create_executive_summary_slide(prs, insights_sections, all_data)
    create_channel_overview_slide(prs, all_data)
    create_content_performance_slide(prs, all_data)
    create_engagement_chart_slide(prs, all_data)
    create_topics_themes_slide(prs, all_data, insights_sections)
    create_posting_frequency_slide(prs, all_data, insights_sections)
    create_gap_analysis_slide(prs, insights_sections)
    create_recommendations_slide(prs, insights_sections, main_company)
    create_rankings_slide(prs, all_data, insights_sections)
    create_summary_slide(prs, all_data, main_company)

    filename = f"video_intelligence_report_{main_company.replace(' ', '_')}.pptx"
    prs.save(filename)
    print(f"PowerPoint saved: {filename}")
    return filename